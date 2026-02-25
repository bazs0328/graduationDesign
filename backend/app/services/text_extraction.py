import logging
import re
import time
from abc import ABC, abstractmethod
from dataclasses import dataclass
from typing import Any, List, Optional, Tuple

import pdfplumber

from app.core.config import settings
from app.services.pdf_layout import (
    ExtractedBlock,
    PageLayoutResult,
    extract_pdf_layout,
)


@dataclass
class ExtractionResult:
    text: str
    page_count: int
    pages: List[str]
    encoding: Optional[str] = None
    blocks: Optional[List[ExtractedBlock]] = None
    page_blocks: Optional[List[PageLayoutResult]] = None
    sidecar: Optional[dict[str, Any]] = None


@dataclass
class OcrPageResult:
    text: str
    engine: str
    avg_confidence: Optional[float] = None
    raw_line_count: int = 0


class OcrEngine(ABC):
    name: str

    @abstractmethod
    def ocr_page(self, image: Any) -> OcrPageResult:
        raise NotImplementedError


FALLBACK_ENCODINGS: Tuple[str, ...] = (
    "utf-8",
    "utf-8-sig",
    "gb18030",
    "gbk",
    "big5",
    "latin-1",
)

logger = logging.getLogger(__name__)

_OCR_ENGINES: dict[str, OcrEngine] = {}
_PREPROCESS_DEPENDENCY_WARNING_LOGGED = False
_PREPROCESS_FAILURE_WARNING_LOGGED = False
_UNKNOWN_OCR_ENGINE_WARNED: set[str] = set()


class RapidOcrEngine(OcrEngine):
    name = "rapidocr"

    def __init__(self) -> None:
        self._engine: Any = None

    def _get_engine(self) -> Any:
        if self._engine is not None:
            return self._engine
        try:
            from rapidocr_onnxruntime import RapidOCR  # type: ignore
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(
                "OCR dependency missing: install rapidocr-onnxruntime."
            ) from exc
        self._engine = RapidOCR()
        return self._engine

    def ocr_page(self, image: Any) -> OcrPageResult:
        try:
            import numpy as np  # type: ignore
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError("OCR dependency missing: install numpy.") from exc

        engine = self._get_engine()
        rgb_image = image.convert("RGB")
        image_array = np.array(rgb_image)

        try:
            ocr_output, _ = engine(image_array)
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError("RapidOCR failed while recognizing the page.") from exc

        lines, confidences = _parse_rapidocr_output(ocr_output)
        text = _normalize_text("\n".join(lines))
        avg_confidence = None
        if confidences:
            avg_confidence = sum(confidences) / len(confidences)
        return OcrPageResult(
            text=text,
            engine=self.name,
            avg_confidence=avg_confidence,
            raw_line_count=len(lines),
        )


class TesseractOcrEngine(OcrEngine):
    name = "tesseract"

    def ocr_page(self, image: Any) -> OcrPageResult:
        try:
            import pytesseract  # type: ignore
            from pytesseract import TesseractNotFoundError  # type: ignore
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(
                "OCR dependency missing: install pytesseract and tesseract-ocr."
            ) from exc

        language = _get_tesseract_language()
        try:
            text = pytesseract.image_to_string(image, lang=language)
        except TesseractNotFoundError as exc:
            raise RuntimeError(
                "Tesseract OCR engine not found. Please install tesseract-ocr."
            ) from exc
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError("Tesseract OCR failed while recognizing the page.") from exc

        normalized_text = _normalize_text(text)
        line_count = len([line for line in normalized_text.splitlines() if line.strip()])
        avg_confidence = _extract_tesseract_avg_confidence(image, language)
        return OcrPageResult(
            text=normalized_text,
            engine=self.name,
            avg_confidence=avg_confidence,
            raw_line_count=line_count,
        )


class CloudOcrEngine(OcrEngine):
    name = "cloud"

    def ocr_page(self, image: Any) -> OcrPageResult:  # noqa: ARG002
        raise RuntimeError("Cloud OCR engine is not implemented yet.")


_OCR_ENGINE_CLASSES: dict[str, type[OcrEngine]] = {
    "rapidocr": RapidOcrEngine,
    "tesseract": TesseractOcrEngine,
    "cloud": CloudOcrEngine,
}


def _normalize_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _detect_encoding(data: bytes) -> Optional[str]:
    try:
        from charset_normalizer import from_bytes  # type: ignore

        result = from_bytes(data).best()
        if result and result.encoding:
            return result.encoding
    except Exception:
        return None
    return None


def _read_text_with_fallback(file_path: str) -> Tuple[str, Optional[str]]:
    with open(file_path, "rb") as f:
        data = f.read()

    encoding = _detect_encoding(data)
    if encoding:
        try:
            return data.decode(encoding), encoding
        except Exception:
            encoding = None

    for enc in FALLBACK_ENCODINGS:
        try:
            return data.decode(enc), enc
        except UnicodeDecodeError:
            continue

    return data.decode("utf-8", errors="ignore"), "utf-8-ignored"


def _safe_int(value: Any, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _safe_float(value: Any) -> Optional[float]:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _parse_csv(value: str | None) -> list[str]:
    if not value:
        return []
    return [part.strip() for part in value.split(",") if part and part.strip()]


def _get_ocr_render_dpi() -> int:
    return max(100, _safe_int(getattr(settings, "ocr_render_dpi", 360), 360))


def _get_ocr_check_pages() -> int:
    return max(1, _safe_int(getattr(settings, "ocr_check_pages", 3), 3))


def _get_ocr_low_confidence_threshold() -> float:
    value = _safe_float(getattr(settings, "ocr_low_confidence_threshold", 0.78))
    if value is None:
        return 0.78
    return min(1.0, max(0.0, value))


def _get_tesseract_language() -> str:
    explicit = (getattr(settings, "ocr_tesseract_language", None) or "").strip()
    if explicit:
        return explicit
    legacy = (getattr(settings, "ocr_language", "chi_sim+eng") or "").strip()
    return legacy or "chi_sim+eng"


def _is_low_text_page(page_text: str, min_text_length: int) -> bool:
    return len(page_text.strip()) < min_text_length


def _is_scanned_pdf(pages: List[str], min_text_length: int) -> bool:
    """Treat PDF as scanned when the first few pages have near-empty text layer."""
    if not pages:
        return True

    sample_count = min(len(pages), _get_ocr_check_pages())
    low_text_pages = sum(1 for page in pages[:sample_count] if _is_low_text_page(page, min_text_length))
    return low_text_pages == sample_count


def _visible_chars(text: str) -> str:
    return "".join(ch for ch in (text or "") if not ch.isspace())


def _page_text_quality_metrics(page_text: str) -> dict[str, float]:
    lines = [line.strip() for line in (page_text or "").splitlines() if line.strip()]
    visible = _visible_chars(page_text)
    visible_len = len(visible)
    if not lines or visible_len == 0:
        return {
            "visible_len": float(visible_len),
            "line_count": float(len(lines)),
            "avg_line_len": 0.0,
            "single_char_line_ratio": 1.0 if lines else 0.0,
            "short_line_ratio": 1.0 if lines else 0.0,
            "symbol_ratio": 1.0 if visible_len else 0.0,
            "latin_ratio": 0.0,
        }

    avg_line_len = sum(len(_visible_chars(line)) for line in lines) / max(1, len(lines))
    single_char_lines = sum(1 for line in lines if len(_visible_chars(line)) <= 1)
    short_lines = sum(1 for line in lines if len(_visible_chars(line)) <= 4)
    symbol_count = sum(
        1
        for ch in visible
        if not (ch.isalnum() or ("\u4e00" <= ch <= "\u9fff"))
    )
    latin_count = sum(1 for ch in visible if ("a" <= ch.lower() <= "z"))
    return {
        "visible_len": float(visible_len),
        "line_count": float(len(lines)),
        "avg_line_len": float(avg_line_len),
        "single_char_line_ratio": single_char_lines / max(1, len(lines)),
        "short_line_ratio": short_lines / max(1, len(lines)),
        "symbol_ratio": symbol_count / max(1, visible_len),
        "latin_ratio": latin_count / max(1, visible_len),
    }


def _score_page_text_quality(page_text: str) -> float:
    m = _page_text_quality_metrics(page_text)
    visible_len = m["visible_len"]
    if visible_len <= 0:
        return 0.0
    length_bonus = min(1.0, visible_len / 300.0) * 0.25
    avg_line_bonus = min(1.0, m["avg_line_len"] / 12.0) * 0.15
    score = 0.4 + length_bonus + avg_line_bonus
    score -= m["single_char_line_ratio"] * 0.35
    score -= m["short_line_ratio"] * 0.2
    score -= m["symbol_ratio"] * 0.2
    # Penalize unusually high latin ratio in mixed/Chinese docs less aggressively.
    if m["latin_ratio"] > 0.75 and m["avg_line_len"] < 5:
        score -= 0.15
    return max(0.0, min(1.0, score))


def _is_garbled_text_page(page_text: str) -> bool:
    if not bool(getattr(settings, "pdf_garbled_ocr_enabled", True)):
        return False
    m = _page_text_quality_metrics(page_text)
    if m["visible_len"] <= 0:
        return False
    single_char_threshold = _safe_float(getattr(settings, "pdf_garbled_single_char_line_ratio", 0.45))
    short_line_threshold = _safe_float(getattr(settings, "pdf_garbled_short_line_ratio", 0.65))
    single_char_threshold = 0.45 if single_char_threshold is None else single_char_threshold
    short_line_threshold = 0.65 if short_line_threshold is None else short_line_threshold
    if m["single_char_line_ratio"] >= max(0.0, single_char_threshold):
        return True
    if m["short_line_ratio"] >= max(0.0, short_line_threshold) and m["avg_line_len"] < 4.0:
        return True
    if m["symbol_ratio"] > 0.35 and m["avg_line_len"] < 5.0:
        return True
    return _score_page_text_quality(page_text) < 0.35


def _extract_rapidocr_bbox_sort_key(row: Any, fallback_index: int) -> tuple[float, float, int]:
    if (
        isinstance(row, (list, tuple))
        and row
        and isinstance(row[0], (list, tuple))
        and row[0]
        and isinstance(row[0][0], (list, tuple))
    ):
        xs: list[float] = []
        ys: list[float] = []
        for pt in row[0]:
            if not isinstance(pt, (list, tuple)) or len(pt) < 2:
                continue
            x = _safe_float(pt[0])
            y = _safe_float(pt[1])
            if x is None or y is None:
                continue
            xs.append(x)
            ys.append(y)
        if xs and ys:
            return min(ys), min(xs), fallback_index
    return float(fallback_index), 0.0, fallback_index


def _parse_rapidocr_output(ocr_output: Any) -> tuple[list[str], list[float]]:
    entries: list[tuple[tuple[float, float, int], str, Optional[float]]] = []
    if not ocr_output:
        return [], []
    if not isinstance(ocr_output, (list, tuple)):
        return [], []

    for row_index, row in enumerate(ocr_output):
        if not isinstance(row, (list, tuple)):
            continue

        text_value: str | None = None
        score_value: Optional[float] = None

        if len(row) >= 3 and isinstance(row[1], str):
            text_value = row[1]
            score_value = _safe_float(row[2])
        elif len(row) >= 2 and isinstance(row[0], str):
            text_value = row[0]
            score_value = _safe_float(row[1])
        else:
            for item in row:
                if text_value is None and isinstance(item, str):
                    text_value = item
                    continue
                if score_value is None:
                    score_value = _safe_float(item)

        if not text_value:
            continue
        normalized_line = _normalize_text(text_value)
        if not normalized_line:
            continue
        if score_value is not None:
            if score_value > 1.0:
                score_value = score_value / 100.0
            score_value = min(1.0, max(0.0, score_value))
        entries.append((_extract_rapidocr_bbox_sort_key(row, row_index), normalized_line, score_value))

    entries.sort(key=lambda item: item[0])
    lines = [item[1] for item in entries]
    confidences = [item[2] for item in entries if item[2] is not None]
    return lines, confidences


def _extract_tesseract_avg_confidence(image: Any, language: str) -> Optional[float]:
    try:
        import pytesseract  # type: ignore
    except Exception:
        return None

    try:
        data = pytesseract.image_to_data(
            image,
            lang=language,
            output_type=pytesseract.Output.DICT,
        )
    except Exception:
        return None

    raw_values = data.get("conf", []) if isinstance(data, dict) else []
    confidences: list[float] = []
    for raw in raw_values:
        score = _safe_float(raw)
        if score is None or score < 0:
            continue
        if score > 1.0:
            score = score / 100.0
        confidences.append(min(1.0, max(0.0, score)))

    if not confidences:
        return None
    return sum(confidences) / len(confidences)


def _get_ocr_engine(name: str) -> OcrEngine:
    normalized = (name or "").strip().lower()
    if normalized not in _OCR_ENGINE_CLASSES:
        raise RuntimeError(f"Unsupported OCR engine: {name}")
    engine = _OCR_ENGINES.get(normalized)
    if engine is None:
        engine = _OCR_ENGINE_CLASSES[normalized]()
        _OCR_ENGINES[normalized] = engine
    return engine


def _get_ocr_engine_chain_names() -> list[str]:
    primary = (getattr(settings, "ocr_engine", "rapidocr") or "rapidocr").strip().lower()
    configured = _parse_csv(getattr(settings, "ocr_fallback_engines", "rapidocr"))
    merged = [primary, *configured]

    result: list[str] = []
    seen: set[str] = set()
    for raw in merged:
        name = (raw or "").strip().lower()
        if not name or name in seen:
            continue
        if name not in _OCR_ENGINE_CLASSES:
            if name not in _UNKNOWN_OCR_ENGINE_WARNED:
                _UNKNOWN_OCR_ENGINE_WARNED.add(name)
                logger.warning("Skip unsupported OCR engine in chain: %s", name)
            continue
        seen.add(name)
        result.append(name)

    if not result:
        result = ["rapidocr"]
    return result


def _should_continue_after_low_confidence(
    result: OcrPageResult,
    *,
    threshold: float,
    has_next_engine: bool,
) -> bool:
    if not has_next_engine:
        return False
    if not result.text:
        return True
    if result.avg_confidence is None:
        return False
    return result.avg_confidence < threshold


def _is_better_ocr_result(candidate: OcrPageResult, current: OcrPageResult) -> bool:
    candidate_conf = candidate.avg_confidence if candidate.avg_confidence is not None else -1.0
    current_conf = current.avg_confidence if current.avg_confidence is not None else -1.0
    if candidate_conf != current_conf:
        return candidate_conf > current_conf
    return len(candidate.text) > len(current.text)


def _run_ocr_with_fallbacks(image: Any, page_num: int) -> OcrPageResult:
    engine_chain = _get_ocr_engine_chain_names()
    confidence_threshold = _get_ocr_low_confidence_threshold()
    best_result: Optional[OcrPageResult] = None
    last_runtime_error: Optional[RuntimeError] = None

    for idx, engine_name in enumerate(engine_chain):
        is_fallback = idx > 0
        has_next_engine = idx < len(engine_chain) - 1
        started = time.perf_counter()
        try:
            result = _get_ocr_engine(engine_name).ocr_page(image)
            elapsed_ms = (time.perf_counter() - started) * 1000
        except RuntimeError as exc:
            elapsed_ms = (time.perf_counter() - started) * 1000
            last_runtime_error = exc
            logger.warning(
                "OCR page=%s engine=%s status=runtime_error fallback=%s elapsed_ms=%.1f error=%s",
                page_num,
                engine_name,
                is_fallback,
                elapsed_ms,
                str(exc),
            )
            continue
        except Exception:  # noqa: BLE001
            elapsed_ms = (time.perf_counter() - started) * 1000
            logger.exception(
                "OCR page=%s engine=%s status=error fallback=%s elapsed_ms=%.1f",
                page_num,
                engine_name,
                is_fallback,
                elapsed_ms,
            )
            continue

        if result.text and (best_result is None or _is_better_ocr_result(result, best_result)):
            best_result = result

        logger.info(
            "OCR page=%s engine=%s status=ok fallback=%s elapsed_ms=%.1f conf=%s lines=%s chars=%s",
            page_num,
            result.engine,
            is_fallback,
            elapsed_ms,
            f"{result.avg_confidence:.3f}" if result.avg_confidence is not None else "n/a",
            result.raw_line_count,
            len(result.text),
        )

        if _should_continue_after_low_confidence(
            result,
            threshold=confidence_threshold,
            has_next_engine=has_next_engine,
        ):
            logger.info(
                "OCR page=%s engine=%s status=low_confidence conf=%s threshold=%.2f action=fallback",
                page_num,
                result.engine,
                f"{result.avg_confidence:.3f}" if result.avg_confidence is not None else "n/a",
                confidence_threshold,
            )
            continue

        if result.text:
            return result

    if best_result is not None:
        return best_result
    if last_runtime_error is not None:
        raise last_runtime_error
    return OcrPageResult(text="", engine=engine_chain[-1], avg_confidence=None, raw_line_count=0)


def _deskew_binary_image(binary_image: Any) -> Any:
    try:
        import cv2  # type: ignore
        import numpy as np  # type: ignore
    except Exception:
        return binary_image

    coords = cv2.findNonZero(255 - binary_image)
    if coords is None or len(coords) < 10:
        return binary_image

    angle = cv2.minAreaRect(coords)[-1]
    if angle < -45:
        angle = -(90 + angle)
    else:
        angle = -angle

    if abs(float(angle)) < 0.2 or abs(float(angle)) > 15:
        return binary_image

    h, w = binary_image.shape[:2]
    matrix = cv2.getRotationMatrix2D((w / 2, h / 2), float(angle), 1.0)
    return cv2.warpAffine(
        binary_image,
        matrix,
        (w, h),
        flags=cv2.INTER_CUBIC,
        borderMode=cv2.BORDER_CONSTANT,
        borderValue=255,
    )


def _preprocess_ocr_image(image: Any) -> Any:
    global _PREPROCESS_DEPENDENCY_WARNING_LOGGED
    global _PREPROCESS_FAILURE_WARNING_LOGGED

    if not bool(getattr(settings, "ocr_preprocess_enabled", True)):
        return image

    try:
        import cv2  # type: ignore
        import numpy as np  # type: ignore
        from PIL import Image  # type: ignore
    except Exception:
        if not _PREPROCESS_DEPENDENCY_WARNING_LOGGED:
            _PREPROCESS_DEPENDENCY_WARNING_LOGGED = True
            logger.warning(
                "OCR preprocessing disabled at runtime because cv2/numpy/Pillow is unavailable"
            )
        return image

    try:
        rgb = image.convert("RGB")
        rgb_array = np.array(rgb)
        gray = cv2.cvtColor(rgb_array, cv2.COLOR_RGB2GRAY)
        denoised = cv2.medianBlur(gray, 3)
        _, binary = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        if bool(getattr(settings, "ocr_deskew_enabled", True)):
            binary = _deskew_binary_image(binary)
        return Image.fromarray(binary)
    except Exception:  # noqa: BLE001
        if not _PREPROCESS_FAILURE_WARNING_LOGGED:
            _PREPROCESS_FAILURE_WARNING_LOGGED = True
            logger.exception("OCR preprocessing failed; falling back to original rendered page image")
        return image


def _render_pdf_page_for_ocr(file_path: str, page_num: int) -> Any:
    try:
        from pdf2image import convert_from_path  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "OCR dependency missing: install pdf2image and poppler-utils."
        ) from exc

    try:
        images = convert_from_path(
            file_path,
            first_page=page_num,
            last_page=page_num,
            dpi=_get_ocr_render_dpi(),
            fmt="png",
        )
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            "Failed to render PDF page for OCR. Please install poppler-utils."
        ) from exc

    if not images:
        return None

    image = images[0]
    processed_image = _preprocess_ocr_image(image)
    if processed_image is not image:
        try:
            image.close()
        except Exception:
            pass
    return processed_image


def _ocr_page(file_path: str, page_num: int) -> OcrPageResult:
    """Run OCR for one PDF page and return normalized text and metadata."""
    image = _render_pdf_page_for_ocr(file_path, page_num)
    if image is None:
        return OcrPageResult(text="", engine="none", avg_confidence=None, raw_line_count=0)

    try:
        return _run_ocr_with_fallbacks(image, page_num)
    finally:
        try:
            image.close()
        except Exception:
            pass


def _merge_page_text_with_ocr(
    original_text: str,
    ocr_text: str,
    *,
    scanned_pdf: bool,
    min_text_length: int,
) -> str:
    if not ocr_text:
        return original_text
    if not original_text:
        return ocr_text
    if scanned_pdf or _is_low_text_page(original_text, min_text_length):
        return ocr_text
    if ocr_text in original_text:
        return original_text
    if original_text in ocr_text:
        return ocr_text
    return f"{original_text}\n{ocr_text}".strip()


def _get_pdf_parser_mode() -> str:
    mode = (getattr(settings, "pdf_parser_mode", "legacy") or "legacy").strip().lower()
    if mode not in {"legacy", "layout", "auto"}:
        return "legacy"
    return mode


def _pick_page_text_after_ocr(
    original_text: str,
    ocr_text: str,
    *,
    scanned_pdf: bool,
    min_text_length: int,
) -> tuple[str, str]:
    """Return (chosen_text, reason)."""
    if not ocr_text:
        return original_text, "no_ocr_text"
    if not original_text:
        return ocr_text, "ocr_empty_text_layer"
    if scanned_pdf or _is_low_text_page(original_text, min_text_length):
        return ocr_text, "ocr_low_text_or_scanned"

    original_score = _score_page_text_quality(original_text)
    ocr_score = _score_page_text_quality(ocr_text)
    if ocr_score >= original_score + 0.05:
        return ocr_text, "ocr_quality_better"
    if original_score >= ocr_score + 0.08:
        return original_text, "text_layer_quality_better"

    if ocr_text in original_text:
        return original_text, "ocr_contained"
    if original_text in ocr_text:
        return ocr_text, "text_layer_contained"
    if len(_visible_chars(ocr_text)) > len(_visible_chars(original_text)) * 1.2:
        return ocr_text, "ocr_much_longer"
    return original_text, "keep_text_layer"


def _apply_pdf_ocr_repair(
    file_path: str,
    *,
    pages: List[str],
    page_blocks: Optional[List[PageLayoutResult]] = None,
) -> tuple[List[str], bool]:
    min_text_length = max(1, settings.ocr_min_text_length)
    scanned_pdf = _is_scanned_pdf(pages, min_text_length)
    if not settings.ocr_enabled:
        if scanned_pdf:
            logger.info("Scanned PDF detected but OCR is disabled")
        return pages, scanned_pdf

    ocr_pages: list[int] = []
    page_reasons: dict[int, list[str]] = {}
    for idx, page_text in enumerate(pages, start=1):
        reasons: list[str] = []
        if scanned_pdf:
            reasons.append("scanned_pdf")
        else:
            if _is_low_text_page(page_text, min_text_length):
                reasons.append("low_text")
            if _is_garbled_text_page(page_text):
                reasons.append("garbled_text")
        if reasons:
            ocr_pages.append(idx)
            page_reasons[idx] = reasons

    for page_num in ocr_pages:
        try:
            ocr_result = _ocr_page(file_path, page_num)
        except RuntimeError:
            if scanned_pdf:
                raise
            logger.exception("Skip OCR on page %s due to setup/runtime error", page_num)
            continue

        if not ocr_result.text:
            continue

        original_page_text = pages[page_num - 1]
        chosen_text, choose_reason = _pick_page_text_after_ocr(
            original_page_text,
            ocr_result.text,
            scanned_pdf=scanned_pdf,
            min_text_length=min_text_length,
        )
        pages[page_num - 1] = chosen_text

        if page_blocks and 0 <= page_num - 1 < len(page_blocks):
            page_block = page_blocks[page_num - 1]
            page_block.text_quality_score = _score_page_text_quality(chosen_text)
            if chosen_text != original_page_text:
                page_block.ocr_override_text = chosen_text

        logger.info(
            "PDF page=%s ocr_trigger=%s choose=%s reason=%s orig_q=%.3f ocr_q=%.3f",
            page_num,
            ",".join(page_reasons.get(page_num, [])) or "none",
            "ocr" if chosen_text == ocr_result.text else "text_layer",
            choose_reason,
            _score_page_text_quality(original_page_text),
            _score_page_text_quality(ocr_result.text),
        )
    return pages, scanned_pdf


def _extract_pdf_legacy(file_path: str) -> ExtractionResult:
    pages: List[str] = []
    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
            pages.append(_normalize_text(page_text))

    pages, _ = _apply_pdf_ocr_repair(file_path, pages=pages, page_blocks=None)

    combined = "\n\n".join([p for p in pages if p])
    return ExtractionResult(text=combined, page_count=page_count, pages=pages)


def _extract_pdf_layout_mode(
    file_path: str,
    *,
    user_id: str | None = None,
    kb_id: str | None = None,
    doc_id: str | None = None,
) -> ExtractionResult:
    parsed = extract_pdf_layout(file_path, user_id=user_id, kb_id=kb_id, doc_id=doc_id)
    pages = [str(p or "") for p in parsed.pages]
    page_blocks = list(parsed.page_blocks)
    pages, _ = _apply_pdf_ocr_repair(file_path, pages=pages, page_blocks=page_blocks)
    combined = "\n\n".join([p for p in pages if p])
    return ExtractionResult(
        text=combined,
        page_count=parsed.page_count,
        pages=pages,
        blocks=list(parsed.blocks),
        page_blocks=page_blocks,
        sidecar=parsed.sidecar,
    )


def _extract_pdf(
    file_path: str,
    *,
    user_id: str | None = None,
    kb_id: str | None = None,
    doc_id: str | None = None,
) -> ExtractionResult:
    mode = _get_pdf_parser_mode()
    if mode in {"layout", "auto"}:
        try:
            return _extract_pdf_layout_mode(
                file_path,
                user_id=user_id,
                kb_id=kb_id,
                doc_id=doc_id,
            )
        except Exception:
            if mode == "layout":
                raise
            logger.exception("PDF layout parser failed, fallback to legacy parser")
    return _extract_pdf_legacy(file_path)


def _extract_text_file(file_path: str) -> ExtractionResult:
    raw_text, encoding = _read_text_with_fallback(file_path)
    normalized = _normalize_text(raw_text)
    return ExtractionResult(text=normalized, page_count=1, pages=[normalized], encoding=encoding)


def _extract_docx(file_path: str) -> ExtractionResult:
    try:
        from docx import Document as WordDocument  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("DOCX dependency missing: install python-docx.") from exc

    try:
        doc = WordDocument(file_path)
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("Failed to read DOCX file.") from exc

    paragraphs: List[str] = []
    for paragraph in doc.paragraphs:
        text = _normalize_text(paragraph.text or "")
        if text:
            paragraphs.append(text)

    combined = "\n\n".join(paragraphs)
    return ExtractionResult(text=combined, page_count=1, pages=[combined])


def _extract_pptx(file_path: str) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("PPTX dependency missing: install python-pptx.") from exc

    try:
        presentation = Presentation(file_path)
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("Failed to read PPTX file.") from exc

    pages: List[str] = []
    for slide in presentation.slides:
        parts: List[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = _normalize_text(getattr(shape, "text", "") or "")
            if shape_text:
                parts.append(shape_text)
        pages.append(_normalize_text("\n\n".join(parts)) if parts else "")

    combined = "\n\n".join([p for p in pages if p])
    return ExtractionResult(text=combined, page_count=len(pages), pages=pages)


def extract_text(
    file_path: str,
    suffix: str,
    *,
    user_id: str | None = None,
    kb_id: str | None = None,
    doc_id: str | None = None,
) -> ExtractionResult:
    normalized_suffix = (suffix or "").lower()
    if normalized_suffix == ".pdf":
        return _extract_pdf(file_path, user_id=user_id, kb_id=kb_id, doc_id=doc_id)
    if normalized_suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    if normalized_suffix == ".docx":
        return _extract_docx(file_path)
    if normalized_suffix == ".pptx":
        return _extract_pptx(file_path)
    raise ValueError(f"Unsupported file type: {suffix}")
