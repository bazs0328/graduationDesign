from pathlib import Path

import pytest

from app.services.text_extraction import extract_text


def test_extract_text_docx_preserves_paragraph_order(tmp_path: Path):
    docx = pytest.importorskip("docx")

    file_path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_paragraph("第一段：线性代数")
    doc.add_paragraph("")
    doc.add_paragraph("第二段：矩阵运算")
    doc.save(file_path)

    result = extract_text(str(file_path), ".docx")

    assert result.page_count == 1
    assert len(result.pages) == 1
    assert result.pages[0] == "第一段：线性代数\n\n第二段：矩阵运算"
    assert result.text == result.pages[0]
    assert result.encoding is None


def test_extract_text_pptx_returns_slide_based_pages(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    file_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]

    slide1 = presentation.slides.add_slide(blank)
    textbox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    text_frame1 = textbox1.text_frame
    text_frame1.text = "第一页标题"
    p1 = text_frame1.add_paragraph()
    p1.text = "第一页内容"

    slide2 = presentation.slides.add_slide(blank)
    textbox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    textbox2.text_frame.text = "第二页内容"

    presentation.slides.add_slide(blank)  # empty slide
    presentation.save(file_path)

    result = extract_text(str(file_path), ".pptx")

    assert result.page_count == 3
    assert len(result.pages) == 3
    assert "第一页标题" in result.pages[0]
    assert "第一页内容" in result.pages[0]
    assert result.pages[1] == "第二页内容"
    assert result.pages[2] == ""
    assert "第一页标题" in result.text
    assert "第二页内容" in result.text


def test_extract_text_docx_empty_content_returns_empty_text(tmp_path: Path):
    docx = pytest.importorskip("docx")

    file_path = tmp_path / "empty.docx"
    doc = docx.Document()
    doc.save(file_path)

    result = extract_text(str(file_path), ".docx")

    assert result.page_count == 1
    assert result.pages == [""]
    assert result.text == ""


def test_extract_text_pptx_empty_slides_return_empty_text(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore

    file_path = tmp_path / "empty.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(file_path)

    result = extract_text(str(file_path), ".pptx")

    assert result.page_count == 1
    assert result.pages == [""]
    assert result.text == ""
